"""Build the Veritype pitch deck. No template required - we draw native PPTX shapes."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

HERE = Path(__file__).parent
SCREENSHOTS = HERE.parent / "docs" / "screenshots"

# ─── Brand palette ───────────────────────────────────────────────────────
AUBERGINE = RGBColor(0x3F, 0x0E, 0x40)
AUBERGINE_DEEP = RGBColor(0x19, 0x17, 0x1D)
INK = RGBColor(0x1D, 0x1C, 0x1D)
INK_2 = RGBColor(0x45, 0x42, 0x45)
INK_3 = RGBColor(0x61, 0x60, 0x61)
HAIRLINE = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF8, 0xF8)
PINK = RGBColor(0xE0, 0x1E, 0x5A)
YELLOW = RGBColor(0xEC, 0xB2, 0x2E)
GREEN = RGBColor(0x2E, 0xB6, 0x7D)
CYAN = RGBColor(0x36, 0xC5, 0xF0)
BLUE = RGBColor(0x12, 0x64, 0xA3)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs, fill=WHITE):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.fill.background()
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Lato", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w
    sh.shadow.inherit = False
    return sh


def add_chip(slide, x, y, text, *, fg=WHITE, bg=AUBERGINE, w=None, h=Inches(0.32), size=10):
    if w is None:
        w = Inches(0.04 * len(text) + 0.4)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    chip.adjustments[0] = 0.5
    chip.fill.solid()
    chip.fill.fore_color.rgb = bg
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Lato"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg


def page_number(slide, n, total):
    add_text(
        slide,
        SLIDE_W - Inches(1.2),
        SLIDE_H - Inches(0.4),
        Inches(1.0),
        Inches(0.3),
        f"{n} / {total}",
        size=10,
        color=INK_3,
        align=PP_ALIGN.RIGHT,
    )


def footer(slide):
    add_text(
        slide,
        Inches(0.5),
        SLIDE_H - Inches(0.4),
        Inches(6),
        Inches(0.3),
        "Veritype : Slack Agent Builder Challenge 2026",
        size=10,
        color=INK_3,
    )


def color_bar(slide, y=0):
    cols = [CYAN, GREEN, YELLOW, PINK]
    seg_w = SLIDE_W / len(cols)
    for i, c in enumerate(cols):
        add_rect(slide, Emu(int(i * seg_w)), Emu(y), Emu(int(seg_w)) + Emu(1), Emu(int(Inches(0.08))), fill=c)


# ─── Slides ──────────────────────────────────────────────────────────────


def slide_cover(prs, total):
    s = blank_slide(prs, fill=AUBERGINE)
    color_bar(s, y=0)

    add_text(
        s,
        Inches(0.7),
        Inches(0.6),
        Inches(6),
        Inches(0.4),
        "SLACK AGENT BUILDER CHALLENGE 2026",
        size=11,
        bold=True,
        color=RGBColor(0xBB, 0xBA, 0xBC),
    )

    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.8), "Veritype", size=110, bold=True, color=WHITE)

    add_text(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(11),
        Inches(0.7),
        "A Slack-native fact-checker that renders chart-rich verdict cards.",
        size=24,
        color=RGBColor(0xDD, 0xDD, 0xDD),
    )

    # datatype expression hint
    add_text(
        s,
        Inches(0.7),
        Inches(4.4),
        Inches(11),
        Inches(0.5),
        "Track 1 : New Slack Agent  ·  TypeScript Coded Agent  ·  MCP + Real-Time Search",
        size=14,
        color=CYAN,
    )

    add_text(
        s,
        Inches(0.7),
        SLIDE_H - Inches(1.0),
        Inches(6),
        Inches(0.5),
        "Tim Dries  ·  Biztory  ·  May 2026",
        size=14,
        color=RGBColor(0xBB, 0xBA, 0xBC),
    )

    add_text(
        s,
        SLIDE_W - Inches(6.7),
        SLIDE_H - Inches(1.0),
        Inches(6),
        Inches(0.5),
        "github.com/timdries/veritype",
        size=14,
        color=CYAN,
        align=PP_ALIGN.RIGHT,
    )


def slide_problem(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)

    add_chip(s, Inches(0.7), Inches(0.7), "01  THE PROBLEM", bg=AUBERGINE, w=Inches(1.9))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(2.0), "Misinformation moves faster\nthan the team can Google.", size=42, bold=True)
    add_text(
        s,
        Inches(0.7),
        Inches(3.0),
        Inches(12),
        Inches(1.1),
        "Stats and half-truths land in channel, get three emoji reacts, and become 'fact' by lunch.\nThe person who pushes back has to do their homework alone, in a side tab, and disrupt the thread.",
        size=18,
        color=INK_2,
    )

    # Quote card
    qx, qy, qw, qh = Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.4)
    add_rect(s, qx, qy, qw, qh, fill=BG, line=HAIRLINE, line_w=Pt(0.75))
    add_rect(s, qx, qy, Emu(int(Inches(0.08))), qh, fill=PINK)

    add_text(s, qx + Inches(0.4), qy + Inches(0.3), qw - Inches(0.8), Inches(0.4), "How it goes today", size=12, bold=True, color=INK_3)

    bullets = [
        "Claim drops in #general at 09:47.",
        "Three reactions, no challenge.",
        "Someone tries to push back at 10:30, has to Google for ten minutes.",
        "Source debate derails the channel for the rest of the morning.",
    ]
    add_text(
        s,
        qx + Inches(0.4),
        qy + Inches(0.8),
        qw - Inches(0.8),
        qh - Inches(1.0),
        "\n".join("·  " + b for b in bullets),
        size=17,
        color=INK,
    )

    footer(s)
    page_number(s, n, total)


def slide_solution(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)
    add_chip(s, Inches(0.7), Inches(0.7), "02  WHAT WE BUILT", bg=AUBERGINE, w=Inches(2.4))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1.5), "Veritype offers a one-tap verification, in thread, in seconds.", size=36, bold=True)

    # Three feature pillars
    pillars = [
        ("Listen", "A Haiku 4.5 classifier reads every channel message and decides what is a verifiable, public-interest claim.", GREEN),
        ("Verify", "Opus 4.7 with web_search fans out, reads primary sources, and produces a JSON verdict with confidence and citations.", YELLOW),
        ("Render", "We paint the verdict as a Datatype-font card. Plain text like {b:20,45,70} renders as a real bar chart.", PINK),
    ]
    pw = Inches(3.9)
    gap = Inches(0.2)
    for i, (title, body, accent) in enumerate(pillars):
        x = Inches(0.7) + (pw + gap) * i
        y = Inches(3.0)
        h = Inches(3.7)
        add_rect(s, x, y, pw, h, fill=WHITE, line=HAIRLINE, line_w=Pt(0.75))
        add_rect(s, x, y, pw, Emu(int(Inches(0.12))), fill=accent)
        add_text(s, x + Inches(0.4), y + Inches(0.5), pw - Inches(0.8), Inches(0.6), title, size=26, bold=True, color=INK)
        add_text(s, x + Inches(0.4), y + Inches(1.3), pw - Inches(0.8), Inches(2.2), body, size=15, color=INK_2)
        # mini code chip
        add_text(
            s,
            x + Inches(0.4),
            y + h - Inches(0.7),
            pw - Inches(0.8),
            Inches(0.4),
            f"step {i+1}",
            size=11,
            color=INK_3,
            bold=True,
        )

    footer(s)
    page_number(s, n, total)


def slide_demo_card(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)
    add_chip(s, Inches(0.7), Inches(0.7), "03  THE VERDICT CARD", bg=AUBERGINE, w=Inches(2.6))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1.5), "Datatype turns plain text into inline charts.", size=36, bold=True)

    # left: image
    img_path = SCREENSHOTS / "sample-mostly_true.png"
    if img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(0.7), Inches(2.8), height=Inches(4.4))

    # right: callouts
    cx = Inches(5.7)
    add_text(s, cx, Inches(2.9), Inches(7.3), Inches(0.5), "Reading the card", size=18, bold=True)
    callouts = [
        ("MOSTLY TRUE pill", "verdict + confidence as {p:98} pie + percentage"),
        ("Source weights", "{b:92,88,74,81,70}: primary docs sit taller than blogs"),
        ("Evidence balance", "three pies for supports / refutes / context"),
        ("Salience sparkline", "{l:22,28,35,30,42,55,70,88}: how often the claim is showing up"),
        ("Citations", "every source has stance, weight, publisher, excerpt"),
    ]
    y = Inches(3.5)
    for label, body in callouts:
        add_text(s, cx, y, Inches(2.4), Inches(0.4), label, size=13, bold=True, color=BLUE)
        add_text(s, cx + Inches(2.4), y, Inches(5.0), Inches(0.6), body, size=13, color=INK_2)
        y += Inches(0.55)

    footer(s)
    page_number(s, n, total)


def slide_architecture(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)
    add_chip(s, Inches(0.7), Inches(0.7), "04  ARCHITECTURE", bg=AUBERGINE, w=Inches(2.3))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1.0), "One brain, three surfaces.", size=36, bold=True)
    add_text(
        s,
        Inches(0.7),
        Inches(2.1),
        Inches(12),
        Inches(0.6),
        "The Slack bot, the MCP server, and the CLI/dashboard all call the same core module.",
        size=16,
        color=INK_2,
    )

    # Top row: 3 surface boxes
    surfaces = [
        ("Slack Bolt app", "Socket Mode, /verify, message events, button handlers", CYAN),
        ("MCP server", "stdio transport, 3 tools, plug into any MCP client", YELLOW),
        ("CLI / Dashboard", "render cards offline, demo + admin UI", PINK),
    ]
    bw, bh = Inches(3.9), Inches(1.4)
    gap = Inches(0.2)
    for i, (t, b, c) in enumerate(surfaces):
        x = Inches(0.7) + (bw + gap) * i
        y = Inches(3.1)
        add_rect(s, x, y, bw, bh, fill=WHITE, line=HAIRLINE, line_w=Pt(1))
        add_rect(s, x, y, Emu(int(Inches(0.08))), bh, fill=c)
        add_text(s, x + Inches(0.35), y + Inches(0.25), bw - Inches(0.5), Inches(0.5), t, size=20, bold=True)
        add_text(s, x + Inches(0.35), y + Inches(0.85), bw - Inches(0.5), Inches(0.6), b, size=12, color=INK_2)

    # core box
    cx, cy, cw, ch = Inches(2.0), Inches(5.0), Inches(9.3), Inches(1.0)
    add_rect(s, cx, cy, cw, ch, fill=AUBERGINE, line=AUBERGINE_DEEP, line_w=Pt(1))
    add_text(s, cx, cy, cw, ch, "src/core : screen → research → render", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # bottom row
    add_text(
        s,
        Inches(0.7),
        Inches(6.3),
        Inches(12),
        Inches(0.6),
        "Haiku 4.5 (screener)  ·  Opus 4.7 + web_search (research)  ·  Puppeteer + Datatype font (PNG card)",
        size=14,
        color=INK_3,
        align=PP_ALIGN.CENTER,
    )

    footer(s)
    page_number(s, n, total)


def slide_demo_stream(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)
    add_chip(s, Inches(0.7), Inches(0.7), "05  IN-CHANNEL UX", bg=AUBERGINE, w=Inches(2.3))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.9), "It feels native because it is native.", size=34, bold=True)
    add_text(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(12),
        Inches(0.5),
        "Ephemeral offer, threaded verdict, Block Kit mrkdwn fallback for accessibility.",
        size=15,
        color=INK_2,
    )
    img = SCREENSHOTS / "0stream.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(2.0), Inches(2.7), height=Inches(4.4))
    footer(s)
    page_number(s, n, total)


def slide_dashboard(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)
    add_chip(s, Inches(0.7), Inches(0.7), "06  WORKSPACE HOME", bg=AUBERGINE, w=Inches(2.5))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.9), "Admins get a Slack-style dashboard.", size=34, bold=True)
    add_text(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(12),
        Inches(0.5),
        "Live KPIs in Datatype charts. Verdict distribution, latency, channel coverage, library.",
        size=15,
        color=INK_2,
    )
    img = SCREENSHOTS / "0dashboard.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(2.0), Inches(2.7), height=Inches(4.4))
    footer(s)
    page_number(s, n, total)


def slide_judges(prs, total, n):
    s = blank_slide(prs, fill=WHITE)
    color_bar(s, y=0)
    add_chip(s, Inches(0.7), Inches(0.7), "07  WHY IT WINS", bg=AUBERGINE, w=Inches(2.2))
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.9), "We score on every judging axis.", size=34, bold=True)

    rows = [
        ("Tech implementation", "Coded Agent in TypeScript. Haiku 4.5 + Opus 4.7 + web_search. MCP server reuse. Puppeteer pipeline.", CYAN),
        ("Design", "Fully Slack-skinned. Aubergine sidebar, workspace rail, Block Kit verdict, Datatype charts.", GREEN),
        ("Idea quality", "Nobody else is rendering charts via a font in Slack. The proactive ephemeral offer is novel UX.", YELLOW),
        ("Impact", "Misinformation friction inside Slack maps directly to trust at every team using it.", PINK),
        ("Required pillars", "Hits MCP + Real-Time Search. Two of three checked off, not just one.", BLUE),
    ]
    y = Inches(2.3)
    for label, body, color in rows:
        add_rect(s, Inches(0.7), y, Emu(int(Inches(0.12))), Inches(0.8), fill=color)
        add_text(s, Inches(1.0), y, Inches(3.5), Inches(0.5), label, size=16, bold=True)
        add_text(s, Inches(4.7), y, Inches(8.0), Inches(0.8), body, size=14, color=INK_2)
        y += Inches(0.95)

    footer(s)
    page_number(s, n, total)


def slide_close(prs, total, n):
    s = blank_slide(prs, fill=AUBERGINE)
    color_bar(s, y=0)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5), "THANK YOU.", size=14, bold=True, color=CYAN)
    add_text(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(12),
        Inches(2.5),
        "Try Veritype yourself.",
        size=72,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.7),
        Inches(4.0),
        Inches(12),
        Inches(0.5),
        "Live workspace : acme-corp.slack.com  ·  Repo : github.com/timdries/veritype",
        size=18,
        color=RGBColor(0xDD, 0xDD, 0xDD),
    )
    add_text(
        s,
        Inches(0.7),
        Inches(4.8),
        Inches(12),
        Inches(0.5),
        "MCP config : drop into Claude Desktop and the same brain is yours.",
        size=18,
        color=RGBColor(0xDD, 0xDD, 0xDD),
    )

    # branded tagline
    add_text(
        s,
        Inches(0.7),
        SLIDE_H - Inches(1.2),
        Inches(12),
        Inches(0.4),
        "Slack-native fact-checking, with charts that ride inside type.",
        size=14,
        color=CYAN,
    )
    page_number(s, n, total)


def main():
    prs = new_deck()
    total = 9

    slide_cover(prs, total)
    slide_problem(prs, total, 2)
    slide_solution(prs, total, 3)
    slide_demo_card(prs, total, 4)
    slide_architecture(prs, total, 5)
    slide_demo_stream(prs, total, 6)
    slide_dashboard(prs, total, 7)
    slide_judges(prs, total, 8)
    slide_close(prs, total, 9)

    out = HERE / "Veritype-pitch-deck.pptx"
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
